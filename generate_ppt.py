from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── 颜色方案 ──
BG_DARK = RGBColor(0x1B, 0x1F, 0x3B)
BG_CARD = RGBColor(0x24, 0x29, 0x4E)
ACCENT = RGBColor(0x6C, 0x63, 0xFF)
ACCENT_LIGHT = RGBColor(0x8B, 0x83, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xB8, 0xD0)
CYAN = RGBColor(0x00, 0xD2, 0xFF)
GREEN = RGBColor(0x00, 0xE6, 0x96)
ORANGE = RGBColor(0xFF, 0x9F, 0x43)
PINK = RGBColor(0xFF, 0x6B, 0x9D)
YELLOW = RGBColor(0xFF, 0xD9, 0x3D)
DIM = RGBColor(0x70, 0x78, 0x90)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if hasattr(shape, 'adjustments') and len(shape.adjustments) > 0:
        shape.adjustments[0] = 0.05
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    p.font.name = "Microsoft YaHei"
    return tf


def add_bullets(slide, left, top, width, height, items, size=15, color=WHITE, bc=ACCENT_LIGHT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        p.space_before = Pt(3)
        rb = p.add_run()
        rb.text = "●  "
        rb.font.size = Pt(9)
        rb.font.color.rgb = bc
        rb.font.name = "Microsoft YaHei"
        rt = p.add_run()
        rt.text = item
        rt.font.size = Pt(size)
        rt.font.color.rgb = color
        rt.font.name = "Microsoft YaHei"
    return tf


def add_bar(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_accent_line(slide, left, top, width):
    return add_bar(slide, left, top, width, Pt(3), ACCENT)


def add_circle_num(slide, left, top, num, color=GREEN):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.32), Inches(0.32))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    c.shadow.inherit = False
    p = c.text_frame.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(11)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Microsoft YaHei"


def page_num(slide, n, total):
    add_text(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.4),
             f"{n}/{total}", size=11, color=DIM, align=PP_ALIGN.RIGHT)


TOTAL = 8

# ============================================================
# Slide 1: 封面 (简洁)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
add_accent_line(s, Inches(4.0), Inches(2.7), Inches(5.3))
add_text(s, Inches(1), Inches(2.9), Inches(11.3), Inches(1),
         "微服务架构的在线书店系统", size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.1), Inches(11.3), Inches(0.5),
         "开题答辩", size=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.0), Inches(11.3), Inches(0.5),
         "答辩人：XXX　　指导教师：XXX　　2026年3月", size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
page_num(s, 1, TOTAL)

# ============================================================
# Slide 2: 技术选型总览
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
add_text(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "01  技术选型与架构方案", size=28, color=WHITE, bold=True)
add_accent_line(s, Inches(0.8), Inches(1.1), Inches(2))

techs = [
    ("语言与框架", ["Go 1.24 + go-micro v4 微服务框架", "Gin HTTP 框架（API 网关）",
                   "Protobuf IDL + gRPC 服务通信"], ACCENT),
    ("服务治理", ["Consul 服务注册/发现/健康检查", "API Gateway：JWT 鉴权 + 令牌桶限流",
                  "8 个微服务按业务领域拆分"], CYAN),
    ("数据存储", ["MySQL 8.0：事务数据(用户/订单/库存/支付)", "MongoDB 7.0：文档数据(图书/会话)",
                  "Redis 7：购物车 + 多级缓存"], GREEN),
    ("搜索/消息/AI", ["Elasticsearch 8：图书全文检索", "RabbitMQ：事件驱动异步解耦",
                      "Milvus 2.4 + OpenAI：AI智能推荐(创新)"], ORANGE),
]

for i, (title, items, color) in enumerate(techs):
    x = Inches(0.8 + i * 3.05)
    add_rect(s, x, Inches(1.5), Inches(2.8), Inches(5.2), BG_CARD)
    add_bar(s, x, Inches(1.5), Inches(2.8), Pt(4), color)
    add_text(s, x + Inches(0.25), Inches(1.8), Inches(2.3), Inches(0.5),
             title, size=18, color=color, bold=True)
    add_bullets(s, x + Inches(0.25), Inches(2.5), Inches(2.3), Inches(3.5),
                items, size=13, color=LIGHT_GRAY, bc=color)

page_num(s, 2, TOTAL)

# ============================================================
# Slide 3: 系统架构图
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
add_text(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "02  系统架构", size=28, color=WHITE, bold=True)
add_accent_line(s, Inches(0.8), Inches(1.1), Inches(2))

# 客户端
add_rect(s, Inches(1.2), Inches(1.4), Inches(10.9), Inches(0.6), BG_CARD)
add_bar(s, Inches(1.2), Inches(1.4), Pt(4), Inches(0.6), ACCENT_LIGHT)
add_text(s, Inches(1.5), Inches(1.46), Inches(1.5), Inches(0.45),
         "客户端", size=14, color=ACCENT_LIGHT, bold=True)
add_text(s, Inches(3.2), Inches(1.46), Inches(8), Inches(0.45),
         "Web / 移动端 / API", size=13, color=LIGHT_GRAY)

add_text(s, Inches(6.2), Inches(2.0), Inches(1.5), Inches(0.3),
         "▼ HTTP REST", size=9, color=DIM, align=PP_ALIGN.CENTER)

# 网关
add_rect(s, Inches(1.2), Inches(2.25), Inches(10.9), Inches(0.6), BG_CARD)
add_bar(s, Inches(1.2), Inches(2.25), Pt(4), Inches(0.6), CYAN)
add_text(s, Inches(1.5), Inches(2.31), Inches(1.5), Inches(0.45),
         "API 网关", size=14, color=CYAN, bold=True)
add_text(s, Inches(3.2), Inches(2.31), Inches(8), Inches(0.45),
         "Gin · JWT 鉴权 · 令牌桶限流 · 路由转发 · CORS", size=13, color=LIGHT_GRAY)

add_text(s, Inches(5.8), Inches(2.85), Inches(2.5), Inches(0.3),
         "▼ gRPC + Protobuf", size=9, color=DIM, align=PP_ALIGN.CENTER)

# 微服务
svc_y = Inches(3.15)
add_text(s, Inches(1.2), svc_y, Inches(6), Inches(0.3),
         "微服务层（go-micro · Consul 服务发现）", size=13, color=GREEN, bold=True)

svcs = [
    ("User", "用户/JWT", "9001", WHITE),
    ("Book", "图书/检索", "9003", WHITE),
    ("Store", "门店/地理", "9002", WHITE),
    ("Inventory", "库存/锁定", "9004", WHITE),
    ("Cart", "购物车", "9005", WHITE),
    ("Order", "订单/事件", "9006", WHITE),
    ("Payment", "支付/退款", "9007", WHITE),
    ("AI", "推荐/对话", "9008", PINK),
]
for i, (name, desc, port, nc) in enumerate(svcs):
    x = Inches(1.2 + i * 1.38)
    add_rect(s, x, svc_y + Inches(0.3), Inches(1.22), Inches(1.15), BG_CARD)
    add_text(s, x + Inches(0.06), svc_y + Inches(0.35), Inches(1.1), Inches(0.4),
             name, size=12, color=nc, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.06), svc_y + Inches(0.72), Inches(1.1), Inches(0.3),
             desc, size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.06), svc_y + Inches(1.05), Inches(1.1), Inches(0.25),
             f":{port}", size=9, color=DIM, align=PP_ALIGN.CENTER)

add_text(s, Inches(5.8), Inches(4.65), Inches(2.5), Inches(0.3),
         "▼ 数据持久化", size=9, color=DIM, align=PP_ALIGN.CENTER)

# 基础设施
inf_y = Inches(4.9)
add_text(s, Inches(1.2), inf_y, Inches(3), Inches(0.3),
         "基础设施层", size=13, color=ORANGE, bold=True)

infras = [
    ("MySQL 8.0", "用户/订单/库存"),
    ("MongoDB 7.0", "图书/会话"),
    ("Redis 7", "购物车/缓存"),
    ("RabbitMQ", "异步消息"),
    ("ES 8", "全文搜索"),
    ("Milvus 2.4", "向量检索"),
    ("Consul", "服务发现"),
]
for i, (name, desc) in enumerate(infras):
    x = Inches(1.2 + i * 1.58)
    add_rect(s, x, inf_y + Inches(0.3), Inches(1.42), Inches(0.85), BG_CARD)
    add_text(s, x + Inches(0.06), inf_y + Inches(0.34), Inches(1.3), Inches(0.35),
             name, size=11, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.06), inf_y + Inches(0.68), Inches(1.3), Inches(0.3),
             desc, size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Docker
add_rect(s, Inches(1.2), Inches(6.25), Inches(10.9), Inches(0.55), BG_CARD)
add_bar(s, Inches(1.2), Inches(6.25), Pt(4), Inches(0.55), PINK)
add_text(s, Inches(1.5), Inches(6.3), Inches(2), Inches(0.4),
         "容器化部署", size=13, color=PINK, bold=True)
add_text(s, Inches(3.5), Inches(6.3), Inches(8), Inches(0.4),
         "Docker Compose 一键编排全部基础设施", size=12, color=LIGHT_GRAY)

page_num(s, 3, TOTAL)

# ============================================================
# Slide 4: 核心业务流程与关键设计
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
add_text(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "03  核心业务流程与关键设计", size=28, color=WHITE, bold=True)
add_accent_line(s, Inches(0.8), Inches(1.1), Inches(2))

# 购书流程
add_text(s, Inches(0.8), Inches(1.35), Inches(5), Inches(0.45),
         "⬥ 购书主流程", size=17, color=CYAN, bold=True)

steps = [
    ("注册/登录", ACCENT), ("搜索图书", CYAN), ("查看详情", GREEN),
    ("加入购物车", ORANGE), ("创建订单", PINK), ("完成支付", YELLOW),
]
for i, (t, c) in enumerate(steps):
    x = Inches(0.9 + i * 1.98)
    add_rect(s, x, Inches(1.8), Inches(1.5), Inches(0.75), BG_CARD)
    add_bar(s, x, Inches(1.8), Inches(1.5), Pt(3), c)
    add_text(s, x + Inches(0.05), Inches(1.92), Inches(1.4), Inches(0.5),
             t, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text(s, x + Inches(1.55), Inches(1.93), Inches(0.4), Inches(0.4),
                 "→", size=16, color=DIM, align=PP_ALIGN.CENTER)

# 左：下单支付
add_rect(s, Inches(0.8), Inches(2.9), Inches(5.6), Inches(4.1), BG_CARD)
add_text(s, Inches(1.1), Inches(3.0), Inches(5), Inches(0.45),
         "⬥ 下单支付流程（事件驱动）", size=16, color=GREEN, bold=True)

osteps = [
    ("1", "用户提交订单", "Order 创建订单记录"),
    ("2", "锁定库存", "Inventory.LockStock 防超卖"),
    ("3", "发起支付", "Payment 处理支付"),
    ("4", "发布支付事件", "Payment → MQ → payment.success"),
    ("5", "确认并扣减", "Order 消费事件 → DeductStock"),
]
for i, (n, t, d) in enumerate(osteps):
    y = Inches(3.55 + i * 0.65)
    add_circle_num(s, Inches(1.15), y, n, GREEN)
    add_text(s, Inches(1.6), y + Inches(0.02), Inches(1.6), Inches(0.3),
             t, size=13, color=WHITE, bold=True)
    add_text(s, Inches(3.3), y + Inches(0.02), Inches(2.8), Inches(0.3),
             d, size=12, color=LIGHT_GRAY)

# 右：关键设计
add_rect(s, Inches(6.9), Inches(2.9), Inches(5.6), Inches(4.1), BG_CARD)
add_text(s, Inches(7.2), Inches(3.0), Inches(5), Inches(0.45),
         "⬥ 关键技术设计", size=16, color=ORANGE, bold=True)

designs = [
    ("库存防超卖", "Lock → Deduct 两阶段，取消时 Release", ACCENT),
    ("事件驱动解耦", "RabbitMQ 异步：支付→订单→库存扣减", CYAN),
    ("混合持久化", "按数据特征选型：MySQL / Mongo / Redis", GREEN),
    ("搜索降级", "ES 全文检索不可用时回退 MongoDB 查询", ORANGE),
    ("地理查询", "MySQL ST_Distance_Sphere 附近门店", PINK),
]
for i, (t, d, c) in enumerate(designs):
    y = Inches(3.55 + i * 0.65)
    add_text(s, Inches(7.3), y, Inches(2), Inches(0.3), t, size=13, color=c, bold=True)
    add_text(s, Inches(9.3), y, Inches(3), Inches(0.3), d, size=12, color=LIGHT_GRAY)

page_num(s, 4, TOTAL)

# ============================================================
# Slide 5: AI 智能模块（创新点）
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
add_text(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "04  创新点：AI 智能推荐模块", size=28, color=WHITE, bold=True)
add_accent_line(s, Inches(0.8), Inches(1.1), Inches(2))

# 左：Multi-Agent
add_rect(s, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4), BG_CARD)
add_text(s, Inches(1.1), Inches(1.65), Inches(5), Inches(0.45),
         "⬥ 多 Agent 协作", size=18, color=CYAN, bold=True)

agents = [
    ("Librarian 图书管理员", "对话式推荐，配备搜索/库存/加购等 6 个工具", ACCENT),
    ("Recommender 推荐师", "RAG + 购买历史 → 结构化推荐列表", GREEN),
    ("Summarizer 摘要师", "生成书籍摘要/主题/目标读者分析", ORANGE),
    ("SmartSearch 搜索师", "自然语言 → 意图解析 → 精确搜索", CYAN),
    ("TasteAnalyzer 偏好分析", "购买历史 → 用户阅读偏好画像", PINK),
]
for i, (name, desc, c) in enumerate(agents):
    y = Inches(2.3 + i * 0.88)
    add_text(s, Inches(1.2), y, Inches(5), Inches(0.32),
             name, size=14, color=c, bold=True)
    add_text(s, Inches(1.2), y + Inches(0.3), Inches(5), Inches(0.4),
             desc, size=12, color=LIGHT_GRAY)

# 右：RAG 流程
add_rect(s, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.4), BG_CARD)
add_text(s, Inches(7.2), Inches(1.65), Inches(5), Inches(0.45),
         "⬥ RAG 检索增强生成", size=18, color=GREEN, bold=True)

rag = [
    ("1", "用户查询", "接收自然语言问题"),
    ("2", "向量化", "text-embedding-3-small 编码"),
    ("3", "语义检索", "Milvus HNSW 近邻搜索"),
    ("4", "库存增强", "批量校验库存，标注可用性"),
    ("5", "上下文注入", "检索结果+会话历史 → Agent"),
    ("6", "生成响应", "GPT-4o 结合上下文回答"),
]
for i, (n, t, d) in enumerate(rag):
    y = Inches(2.35 + i * 0.75)
    add_circle_num(s, Inches(7.3), y, n, GREEN)
    add_text(s, Inches(7.8), y - Inches(0.02), Inches(1.8), Inches(0.32),
             t, size=14, color=WHITE, bold=True)
    add_text(s, Inches(9.5), y - Inches(0.02), Inches(2.8), Inches(0.32),
             d, size=12, color=LIGHT_GRAY)

page_num(s, 5, TOTAL)

# ============================================================
# Slide 6: 当前进展
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
add_text(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "05  当前进展", size=28, color=WHITE, bold=True)
add_accent_line(s, Inches(0.8), Inches(1.1), Inches(2))

AMBER = RGBColor(0xFF, 0xC1, 0x07)

modules = [
    ("用户服务", "注册、登录、JWT 鉴权、个人资料、地址管理", "初步开发", AMBER),
    ("图书服务", "MongoDB 存储、ES 全文搜索、分类、详情", "初步开发", AMBER),
    ("门店服务", "门店 CRUD、地理查询（附近门店/半径搜索）", "初步开发", AMBER),
    ("库存服务", "库存查询、Lock/Deduct/Release 机制", "初步开发", AMBER),
    ("购物车服务", "Redis 购物车、库存校验", "初步开发", AMBER),
    ("订单服务", "创建/取消/状态流转、MQ 消费支付事件", "初步开发", AMBER),
    ("支付服务", "支付处理、事务控制、MQ 发布事件", "初步开发", AMBER),
    ("AI 服务", "多 Agent、RAG、向量搜索、对话推荐", "初步开发", AMBER),
    ("API 网关", "Gin 路由、JWT 中间件、限流", "初步开发", AMBER),
    ("基础设施", "Docker Compose 编排", "初步搭建", AMBER),
]

hy = Inches(1.5)
add_rect(s, Inches(0.8), hy, Inches(11.7), Inches(0.55), ACCENT)
add_text(s, Inches(1.0), hy + Inches(0.08), Inches(2), Inches(0.4),
         "模块", size=14, color=WHITE, bold=True)
add_text(s, Inches(3.2), hy + Inches(0.08), Inches(6), Inches(0.4),
         "已实现功能（Demo 阶段）", size=14, color=WHITE, bold=True)
add_text(s, Inches(10.2), hy + Inches(0.08), Inches(2.2), Inches(0.4),
         "当前状态", size=14, color=WHITE, bold=True)

for i, (mod, desc, status, sc) in enumerate(modules):
    y = Inches(2.15 + i * 0.5)
    bg = BG_CARD if i % 2 == 0 else BG_DARK
    add_rect(s, Inches(0.8), y, Inches(11.7), Inches(0.45), bg)
    add_bar(s, Inches(0.8), y, Pt(3), Inches(0.45), AMBER)
    add_text(s, Inches(1.0), y + Inches(0.06), Inches(2), Inches(0.35),
             mod, size=13, color=WHITE, bold=True)
    add_text(s, Inches(3.2), y + Inches(0.06), Inches(6.5), Inches(0.35),
             desc, size=12, color=LIGHT_GRAY)
    add_text(s, Inches(10.2), y + Inches(0.06), Inches(2.2), Inches(0.35),
             status, size=12, color=sc)

add_text(s, Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.35),
         "* 以上模块均已完成初步开发，可运行基本流程，但尚未经过完整测试、性能验证和生产部署",
         size=11, color=DIM)

page_num(s, 6, TOTAL)

# ============================================================
# Slide 7: 待完成部分与后续计划
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)
add_text(s, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "06  待完成工作与后续计划", size=28, color=WHITE, bold=True)
add_accent_line(s, Inches(0.8), Inches(1.1), Inches(2))

# 左：待完成
add_rect(s, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4), BG_CARD)
add_text(s, Inches(1.1), Inches(1.65), Inches(5), Inches(0.45),
         "⬥ 待完成工作", size=18, color=ORANGE, bold=True)

todos = [
    "各服务功能完善与边界场景补全",
    "前端页面开发（Web 端用户界面）",
    "服务间联调与端到端集成测试",
    "单元测试与接口测试覆盖",
    "性能压测与瓶颈优化（高并发场景）",
    "系统安全加固（参数校验、注入防护等）",
    "生产环境部署与运维方案",
    "毕业论文撰写",
]
add_bullets(s, Inches(1.1), Inches(2.3), Inches(5), Inches(4.5),
            todos, size=14, color=LIGHT_GRAY, bc=ORANGE)

# 右：后续计划
add_rect(s, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.4), BG_CARD)
add_text(s, Inches(7.2), Inches(1.65), Inches(5), Inches(0.45),
         "⬥ 后续计划", size=18, color=CYAN, bold=True)

plans = [
    ("3月", "功能完善，补全各服务边界逻辑", ACCENT),
    ("3~4月", "前端开发与服务联调", CYAN),
    ("4月", "单元测试、集成测试", GREEN),
    ("4~5月", "性能压测、安全加固、部署上线", ORANGE),
    ("5月", "毕业论文初稿撰写", PINK),
    ("5~6月", "论文修改完善与答辩准备", ACCENT_LIGHT),
]
for i, (time, task, c) in enumerate(plans):
    y = Inches(2.35 + i * 0.78)
    add_circle_num(s, Inches(7.3), y, str(i + 1), c)
    add_text(s, Inches(7.8), y - Inches(0.02), Inches(1.3), Inches(0.35),
             time, size=14, color=c, bold=True)
    add_text(s, Inches(9.2), y - Inches(0.02), Inches(3), Inches(0.35),
             task, size=14, color=LIGHT_GRAY)

page_num(s, 7, TOTAL)

# ============================================================
# Slide 8: 谢谢 / Q&A
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s, BG_DARK)

add_accent_line(s, Inches(5.2), Inches(2.8), Inches(2.9))
add_text(s, Inches(1), Inches(3.0), Inches(11.3), Inches(1),
         "谢谢各位老师", size=38, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.3), Inches(11.3), Inches(0.6),
         "请批评指正！", size=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

page_num(s, 8, TOTAL)

# ── 保存 ──
out = "/Users/qiwang/code/go/book-e-commerce-micro/微服务架构的在线书店系统_开题答辩.pptx"
prs.save(out)
print(f"PPT 已生成: {out}")
